import os
import tempfile
from typing import List, Dict, Any
import PyPDF2
from docx import Document
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from a PDF file."""
    text = ""
    try:
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page_num in range(len(reader.pages)):
                page = reader.pages[page_num]
                text += page.extract_text() + "\n\n"
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
    return text

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from a DOCX file."""
    text = ""
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
    return text

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from a PPTX file."""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
    return text

def extract_text_from_file(file_path: str) -> str:
    """
    Extract text from a file based on its extension.
    Supports PDF, DOCX, and PPTX files.
    """
    file_extension = os.path.splitext(file_path)[1].lower()
    
    if file_extension == '.pdf':
        return extract_text_from_pdf(file_path)
    elif file_extension == '.docx':
        return extract_text_from_docx(file_path)
    elif file_extension == '.pptx' or file_extension == '.ppt':
        return extract_text_from_pptx(file_path)
    else:
        return ""

def split_text(text: str, filename: str) -> List[Dict[str, Any]]:
    """
    Split text into chunks and attach metadata.
    
    Args:
        text: The text to split
        filename: Original filename to include in metadata
        
    Returns:
        List of dictionaries with text and metadata
    """
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=50,
        separators=["\n\n", "\n", ".", " ", ""]
    )
    
    chunks = text_splitter.split_text(text)
    
    # Create documents with metadata
    docs = []
    for i, chunk in enumerate(chunks):
        docs.append({
            "text": chunk,
            "metadata": {
                "source": filename,
                "chunk": i
            }
        })
    
    return docs